# -*- coding: utf-8 -*-
"""Generate a PCI DSS training deck with editable architecture-style diagrams."""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent / ".vendor"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("PCI_DSS_培训要点清单_架构图版_v3.pptx")
FONT = "Microsoft YaHei"

COLORS = {
    "bg": "F4F7FA",
    "ink": "17212B",
    "muted": "64748B",
    "line": "B8C2CC",
    "panel": "FFFFFF",
    "panel2": "EAF0F6",
    "navy": "14324A",
    "teal": "0F766E",
    "cyan": "0891B2",
    "green": "4D7C0F",
    "orange": "C2410C",
    "red": "B91C1C",
    "yellow": "CA8A04",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_text(slide, text, x, y, w, h, size=18, color="ink", bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, fill=None, line=None,
             margin=0.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 2)
    tf.margin_bottom = Inches(margin / 2)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(COLORS[fill])
    else:
        box.fill.background()
    if line:
        box.line.color.rgb = rgb(COLORS[line])
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box


def add_box(slide, text, x, y, w, h, fill="panel", line="line", color="ink",
            size=13, bold=False, radius=True, align=PP_ALIGN.CENTER):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.color.rgb = rgb(COLORS[line])
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(COLORS[color])
    return shp


def add_arrow(slide, x, y, w, h, fill="teal", text="", size=10):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.fill.background()
    if text:
        tf = shp.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = rgb("FFFFFF")
    return shp


def add_down_arrow(slide, x, y, w, h, fill="teal", text=""):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.fill.background()
    if text:
        shp.text_frame.text = text
    return shp


def add_zone(slide, label, x, y, w, h, fill="panel2", line="teal", dashed=False):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.color.rgb = rgb(COLORS[line])
    shp.line.width = Pt(1.5)
    if dashed:
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_text(slide, label, x + 0.12, y + 0.05, w - 0.24, 0.3, size=10, color=line, bold=True)
    return shp


def add_header(slide, title, kicker=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
    if kicker:
        add_text(slide, kicker, 0.55, 0.25, 5.4, 0.28, size=10, color="teal", bold=True)
    add_text(slide, title, 0.55, 0.52, 7.4, 0.55, size=28, color="ink", bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.16), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS["line"])
    line.line.fill.background()


def add_bullets(slide, bullets, x=0.65, y=1.45, w=5.25, h=4.8, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS["ink"])
    return box


def add_small_caption(slide, text, x, y, w, color="muted"):
    return add_text(slide, text, x, y, w, 0.28, size=9, color=color, align=PP_ALIGN.CENTER)


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r_idx == 0 or c_idx == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = rgb(COLORS["ink"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLORS["panel2"] if r_idx == 0 else "FFFFFF")
    return table_shape


def draw_title_visual(slide):
    add_zone(slide, "PCI DSS 决策地图", 6.35, 1.45, 5.95, 4.85, fill="panel", line="navy")
    nodes = [
        ("1 合规判断\nLevel / SAQ / Scope", 6.75, 2.0, "teal"),
        ("2 产品落地\nToken / PSP / CIT-MIT", 9.1, 2.0, "cyan"),
        ("3 技术架构\nVault / Key / Frontend", 6.75, 4.05, "orange"),
        ("4 运营审计\nTest / Access / IR", 9.1, 4.05, "green"),
    ]
    for text, x, y, color in nodes:
        add_box(slide, text, x, y, 1.95, 1.05, fill="panel2", line=color, color="ink", size=11, bold=True)
    add_arrow(slide, 8.55, 2.34, 0.45, 0.27, fill="teal")
    add_down_arrow(slide, 9.77, 3.12, 0.35, 0.75, fill="cyan")
    add_arrow(slide, 8.55, 4.39, 0.45, 0.27, fill="orange")
    add_down_arrow(slide, 7.43, 3.12, 0.35, 0.75, fill="green")
    add_box(slide, "目标：减少卡数据暴露面\n把合规成本转成业务收益", 7.55, 5.42, 3.55, 0.62, fill="navy", line="navy", color="panel", size=12, bold=True)


def draw_toc_visual(slide):
    labels = [("判断", "Scope 是否进来", "teal"), ("产品", "怎样少碰卡", "cyan"), ("架构", "怎样隔离和加密", "orange"), ("运营", "怎样持续证明", "green")]
    for i, (title, subtitle, color) in enumerate(labels):
        x = 6.35 + i * 1.45
        add_box(slide, title, x, 1.7, 1.1, 0.6, fill=color, line=color, color="panel", size=15, bold=True)
        add_box(slide, subtitle, x - 0.08, 2.52, 1.26, 1.28, fill="panel", line=color, size=10)
        add_down_arrow(slide, x + 0.35, 2.28, 0.34, 0.22, fill=color)
    add_arrow(slide, 6.5, 4.35, 5.05, 0.42, fill="navy", text="要求是什么 → 怎么降成本 → 怎么落地 → 怎么证明", size=11)
    add_box(slide, "持续合规运营", 7.85, 5.25, 2.3, 0.62, fill="panel2", line="navy", size=13, bold=True)


def draw_module_visual(slide, module_no, labels):
    add_zone(slide, f"模块 {module_no} 讲解路径", 6.25, 1.45, 5.9, 4.85, fill="panel", line="navy")
    y = 1.95
    colors = ["teal", "cyan", "orange", "green"]
    for idx, label in enumerate(labels):
        add_box(slide, label, 6.75, y + idx * 1.05, 4.55, 0.62, fill="panel2", line=colors[idx % 4], size=13, bold=True)
        if idx < len(labels) - 1:
            add_down_arrow(slide, 8.85, y + idx * 1.05 + 0.66, 0.38, 0.32, fill=colors[idx % 4])
    add_box(slide, "用架构边界解释 PCI 责任", 7.45, 5.65, 3.25, 0.45, fill="navy", line="navy", color="panel", size=11, bold=True)


def draw_boundary_visual(slide):
    add_zone(slide, "责任边界 / Data Flow", 6.15, 1.42, 6.25, 4.95, fill="panel", line="navy")
    add_box(slide, "买家浏览器", 6.55, 2.05, 1.35, 0.55, fill="panel2", line="teal", size=11, bold=True)
    add_box(slide, "商户系统\nAPI / DB / Logs", 8.15, 1.8, 1.62, 1.05, fill="panel2", line="orange", size=10, bold=True)
    add_box(slide, "PSP / 收单行", 10.1, 2.05, 1.35, 0.55, fill="panel2", line="cyan", size=11, bold=True)
    add_arrow(slide, 7.88, 2.22, 0.35, 0.22, fill="teal")
    add_arrow(slide, 9.74, 2.22, 0.35, 0.22, fill="cyan")
    add_zone(slide, "CDE Scope", 7.82, 3.28, 2.45, 1.25, fill="panel2", line="red", dashed=True)
    add_box(slide, "Vault", 8.05, 3.72, 0.72, 0.42, fill="panel", line="red", size=9, bold=True)
    add_box(slide, "密钥", 8.98, 3.72, 0.72, 0.42, fill="panel", line="red", size=9, bold=True)
    add_box(slide, "PCI v4.0\n持续监控", 7.05, 5.35, 4.3, 0.52, fill="navy", line="navy", color="panel", size=12, bold=True)


def draw_license_visual(slide):
    img = Path(__file__).resolve().parent / "assets" / "pci_vs_payment_license.png"
    add_zone(slide, "参考图：监管资质 vs 安全合规", 6.05, 1.34, 6.35, 5.45, fill="panel", line="navy")
    slide.shapes.add_picture(str(img), Inches(6.25), Inches(1.58), height=Inches(4.95))
    add_box(slide, "讲解重点：PCI 是接入卡组网络的安全底线；支付牌照是本地经营许可，两者不是同一件事。", 6.52, 6.22, 5.42, 0.42, fill="panel2", line="navy", size=9, bold=True)


def draw_full_image(slide, asset_name):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")
    img = Path(__file__).resolve().parent / "assets" / asset_name
    slide.shapes.add_picture(str(img), Inches(1.9), Inches(0.15), height=Inches(7.05))


def draw_case_card(slide, case):
    add_zone(slide, "短案例：把概念落到决策", 0.65, 1.35, 12.05, 5.25, fill="panel", line="navy")
    cards = [
        ("场景", case["scene"], "teal", 0.95, 2.0),
        ("误判", case["mistake"], "orange", 6.95, 2.0),
        ("影响", case["impact"], "red", 0.95, 4.05),
        ("正确做法", case["fix"], "green", 6.95, 4.05),
    ]
    for title, body, color, x, y in cards:
        add_box(slide, title, x, y, 1.1, 0.42, fill=color, line=color, color="panel", size=11, bold=True)
        add_box(slide, body, x, y + 0.55, 5.15, 1.15, fill="panel2", line=color, size=14, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide, case["takeaway"], 2.2, 6.15, 8.9, 0.42, fill="navy", line="navy", color="panel", size=11, bold=True)


def draw_verification_visual(slide):
    add_zone(slide, "认证判断链", 6.15, 1.45, 6.25, 4.9, fill="panel", line="navy")
    nodes = [("交易量", "Level 1-4"), ("身份", "Merchant / SP"), ("接入方式", "SAQ A / A-EP / C / D"), ("材料", "AOC / ROC / ASV")]
    x = 6.45
    for i, (a, b) in enumerate(nodes):
        add_box(slide, f"{a}\n{b}", x + i * 1.35, 2.15, 1.12, 1.0, fill="panel2", line=["teal", "cyan", "orange", "green"][i], size=9, bold=True)
        if i < 3:
            add_arrow(slide, x + i * 1.35 + 1.12, 2.5, 0.28, 0.18, fill="navy")
    add_box(slide, "产品形态决定 SAQ\n安全团队只能验证事实", 7.1, 4.35, 4.45, 0.72, fill="navy", line="navy", color="panel", size=12, bold=True)
    add_small_caption(slide, "Hosted Page → A    Elements → A-EP    Direct API → D", 6.7, 5.45, 5.0)


def draw_saq_tree_visual(slide):
    add_zone(slide, "SAQ 决策树", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "卡数据是否进入\n商户系统？", 8.25, 1.85, 1.75, 0.72, fill="navy", line="navy", color="panel", size=11, bold=True)
    add_down_arrow(slide, 8.9, 2.62, 0.38, 0.35, fill="navy")
    add_box(slide, "否：Hosted Page", 6.65, 3.05, 1.55, 0.55, fill="panel2", line="teal", size=10, bold=True)
    add_box(slide, "部分：iframe / Elements", 8.55, 3.05, 1.55, 0.55, fill="panel2", line="cyan", size=9, bold=True)
    add_box(slide, "是：Direct API", 10.45, 3.05, 1.25, 0.55, fill="panel2", line="red", size=9, bold=True)
    add_box(slide, "SAQ A\n商户最小 Scope", 6.65, 4.25, 1.55, 0.7, fill="panel", line="teal", size=10, bold=True)
    add_box(slide, "SAQ A-EP\n前端支付页担责", 8.55, 4.25, 1.55, 0.7, fill="panel", line="cyan", size=9, bold=True)
    add_box(slide, "SAQ D\n完整 CDE 控制", 10.45, 4.25, 1.25, 0.7, fill="panel", line="red", size=9, bold=True)
    add_down_arrow(slide, 7.22, 3.65, 0.32, 0.42, fill="teal")
    add_down_arrow(slide, 9.15, 3.65, 0.32, 0.42, fill="cyan")
    add_down_arrow(slide, 10.9, 3.65, 0.32, 0.42, fill="red")
    add_box(slide, "SAQ 不是安全团队选择题，而是产品接入形态的结果", 6.85, 5.55, 4.85, 0.42, fill="panel2", line="navy", size=10, bold=True)


def draw_scope_visual(slide):
    add_zone(slide, "CDE 与旁路系统", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_zone(slide, "CDE", 7.45, 2.0, 2.45, 2.15, fill="panel2", line="red", dashed=True)
    for label, x, y in [("支付 API", 7.75, 2.48), ("Vault", 8.75, 2.48), ("DB", 7.75, 3.22), ("跳板机", 8.75, 3.22)]:
        add_box(slide, label, x, y, 0.78, 0.42, fill="panel", line="red", size=8, bold=True)
    outer = [("ELK", 6.55, 1.9), ("APM", 10.45, 1.9), ("客服", 6.55, 4.6), ("BI", 10.45, 4.6), ("CI/CD", 8.35, 5.15)]
    for label, x, y in outer:
        add_box(slide, label, x, y, 0.9, 0.42, fill="panel2", line="orange", size=9, bold=True)
    add_arrow(slide, 7.25, 2.05, 0.34, 0.18, fill="orange")
    add_arrow(slide, 9.82, 2.05, 0.34, 0.18, fill="orange")
    add_box(slide, "一条完整 PAN 日志\n会把整套日志系统拉进 Scope", 7.05, 4.55, 3.75, 0.52, fill="red", line="red", color="panel", size=10, bold=True)


def draw_cde_zoom_visual(slide):
    add_zone(slide, "CDE Scope 放大图", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_zone(slide, "CDE 核心区", 7.45, 2.0, 2.35, 2.15, fill="panel2", line="red", dashed=True)
    core = [("支付 API", 7.72, 2.45), ("Vault", 8.62, 2.45), ("PAN DB", 7.72, 3.22), ("KMS/HSM", 8.62, 3.22)]
    for label, x, y in core:
        add_box(slide, label, x, y, 0.78, 0.42, fill="panel", line="red", size=8, bold=True)
    satellites = [
        ("日志 ELK", 6.45, 1.85, "orange"), ("APM Trace", 10.45, 1.85, "orange"),
        ("客服后台", 6.45, 4.2, "cyan"), ("BI / 数仓", 10.45, 4.2, "cyan"),
        ("MQ / Cache", 8.0, 5.05, "green"), ("CI/CD", 9.3, 5.05, "green"),
    ]
    for label, x, y, color in satellites:
        add_box(slide, label, x, y, 1.0, 0.42, fill="panel", line=color, size=8, bold=True)
    add_arrow(slide, 7.25, 2.04, 0.35, 0.16, fill="orange")
    add_arrow(slide, 9.78, 2.04, 0.35, 0.16, fill="orange")
    add_arrow(slide, 7.25, 4.37, 0.35, 0.16, fill="cyan")
    add_arrow(slide, 9.78, 4.37, 0.35, 0.16, fill="cyan")
    add_box(slide, "判断标准：是否存储、处理、传输卡数据，或能影响 CDE 安全", 6.7, 5.65, 5.05, 0.42, fill="navy", line="navy", color="panel", size=10, bold=True)


def draw_three_lanes_visual(slide):
    add_zone(slide, "资金流 / 数据流 / 责任流", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    lanes = [("资金流", 1.85, "green"), ("卡数据流", 3.25, "red"), ("合规责任流", 4.65, "orange")]
    for label, y, color in lanes:
        add_box(slide, label, 6.45, y, 0.95, 0.48, fill=color, line=color, color="panel", size=9, bold=True)
        add_box(slide, "商户", 7.65, y, 0.8, 0.48, fill="panel2", line=color, size=8, bold=True)
        add_arrow(slide, 8.45, y + 0.14, 0.38, 0.18, fill=color)
        add_box(slide, "PSP / 收单", 8.92, y, 1.0, 0.48, fill="panel2", line=color, size=8, bold=True)
        add_arrow(slide, 9.92, y + 0.14, 0.38, 0.18, fill=color)
        add_box(slide, "卡组 / 发卡行", 10.42, y, 1.15, 0.48, fill="panel2", line=color, size=8, bold=True)
    add_box(slide, "同一笔交易里，钱、卡数据、合规责任可能走不同路径", 6.75, 5.65, 4.85, 0.42, fill="navy", line="navy", color="panel", size=10, bold=True)


def draw_product_visual(slide):
    add_zone(slide, "接入形态与 Scope", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    rows = [("Hosted Page", "SAQ A", "最小", "teal"), ("Elements", "SAQ A-EP", "前端担责", "cyan"), ("Direct API", "SAQ D", "最大", "red")]
    for i, (name, saq, scope, color) in enumerate(rows):
        y = 1.82 + i * 1.1
        add_box(slide, name, 6.55, y, 1.42, 0.52, fill="panel2", line=color, size=10, bold=True)
        add_arrow(slide, 8.03, y + 0.15, 0.55, 0.2, fill=color)
        add_box(slide, saq, 8.7, y, 1.05, 0.52, fill="panel", line=color, size=10, bold=True)
        add_box(slide, scope, 10.0, y, 1.3, 0.52, fill="panel", line=color, size=10)
    add_box(slide, "PayFac：平台扛 Level 1\nClick to Pay / Passkey：卡组织或设备托管", 6.85, 5.2, 4.75, 0.75, fill="navy", line="navy", color="panel", size=11, bold=True)


def draw_citmit_visual(slide):
    add_zone(slide, "CIT / MIT 生命周期", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "CIT\n买家在场\n3DS / SCA", 6.55, 2.0, 1.45, 0.9, fill="panel2", line="teal", size=10, bold=True)
    add_arrow(slide, 8.05, 2.32, 0.6, 0.22, fill="teal")
    add_box(slide, "授权凭证\nToken / Agreement", 8.75, 2.0, 1.45, 0.9, fill="panel2", line="cyan", size=10, bold=True)
    add_arrow(slide, 10.25, 2.32, 0.6, 0.22, fill="cyan")
    add_box(slide, "MIT\n订阅续费\n重试 / 过期更新", 10.95, 2.0, 1.0, 0.9, fill="panel2", line="orange", size=9, bold=True)
    add_box(slide, "首次建立授权，后续不能超出授权范围", 7.0, 4.25, 4.65, 0.55, fill="navy", line="navy", color="panel", size=12, bold=True)
    add_small_caption(slide, "订阅业务重点：扣款失败重试 + Network Token 卡过期自动更新", 6.55, 5.35, 5.35)


def draw_citmit_sequence_visual(slide):
    add_zone(slide, "3DS + CIT/MIT 时序", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    actors = [("用户", 6.55), ("商户", 7.75), ("PSP", 8.95), ("3DS / 发卡行", 10.15), ("Token / Vault", 11.35)]
    for label, x in actors:
        add_box(slide, label, x, 1.75, 0.92, 0.4, fill="panel2", line="navy", size=8, bold=True)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.45), Inches(2.25), Inches(0.01), Inches(3.0))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(COLORS["line"])
        line.line.fill.background()
    steps = [
        ("CIT 首笔", 2.45, 6.95, 8.95, "teal"),
        ("3DS/SCA", 2.95, 8.95, 10.15, "orange"),
        ("返回授权凭证", 3.45, 10.15, 8.95, "orange"),
        ("保存 token / agreement", 3.95, 8.95, 11.35, "cyan"),
        ("MIT 续费扣款", 4.65, 7.75, 8.95, "green"),
    ]
    for text, y, x1, x2, color in steps:
        add_arrow(slide, min(x1, x2) + 0.55, y, abs(x2 - x1) - 0.1, 0.18, fill=color, text=text, size=7)
    add_box(slide, "首笔必须建立授权；后续 MIT 必须引用首次凭证", 6.95, 5.65, 4.55, 0.42, fill="navy", line="navy", color="panel", size=10, bold=True)


def draw_token_visual(slide):
    add_zone(slide, "Token 编排", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "用户支付方式", 6.55, 3.0, 1.25, 0.55, fill="panel2", line="teal", size=10, bold=True)
    add_arrow(slide, 7.82, 3.16, 0.42, 0.2, fill="teal")
    add_box(slide, "Token Layer", 8.35, 2.75, 1.5, 1.05, fill="navy", line="navy", color="panel", size=12, bold=True)
    targets = [("PSP token\n换 PSP 失效", 10.15, 1.75, "orange"), ("Merchant token\n跨 PSP", 10.15, 3.0, "cyan"), ("Network token\n授权率 +2~5%", 10.15, 4.25, "green")]
    for text, x, y, color in targets:
        add_arrow(slide, 9.82, y + 0.28, 0.35, 0.18, fill=color)
        add_box(slide, text, x, y, 1.55, 0.7, fill="panel2", line=color, size=8, bold=True)
    add_box(slide, "多 PSP 策略会倒逼独立 Vault 或 Network Token", 6.85, 5.45, 4.75, 0.42, fill="panel", line="navy", size=10, bold=True)


def draw_network_token_lifecycle_visual(slide):
    add_zone(slide, "Network Token 生命周期", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    steps = [
        ("1 请求 token", 6.55, 2.15, "teal"),
        ("2 卡组发放\nNetwork Token", 8.0, 2.15, "cyan"),
        ("3 授权交易", 9.65, 2.15, "green"),
        ("4 卡过期/换卡\n自动更新", 8.0, 3.75, "orange"),
        ("5 多 PSP 路由\n仍用同一凭证体系", 9.65, 3.75, "navy"),
    ]
    for text, x, y, color in steps:
        add_box(slide, text, x, y, 1.25, 0.68, fill="panel2", line=color, size=8, bold=True)
    add_arrow(slide, 7.78, 2.38, 0.35, 0.18, fill="teal")
    add_arrow(slide, 9.25, 2.38, 0.35, 0.18, fill="cyan")
    add_down_arrow(slide, 10.05, 2.9, 0.28, 0.58, fill="green")
    add_arrow(slide, 9.25, 3.98, 0.35, 0.18, fill="orange")
    add_box(slide, "价值：少碰 PAN、卡生命周期自动更新、授权率提升、跨 PSP 更稳", 6.55, 5.45, 5.35, 0.42, fill="panel", line="navy", size=10, bold=True)


def draw_vault_visual(slide):
    add_zone(slide, "Vault 编排架构", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "商户支付服务", 6.55, 3.0, 1.25, 0.55, fill="panel2", line="teal", size=10, bold=True)
    add_arrow(slide, 7.82, 3.16, 0.42, 0.2, fill="teal")
    add_box(slide, "独立 Vault\nToken Map / PAN", 8.35, 2.65, 1.5, 1.25, fill="navy", line="navy", color="panel", size=11, bold=True)
    for i, (label, color) in enumerate([("PSP A\n低成本", "cyan"), ("PSP B\n高授权", "green"), ("PSP C\n容灾", "orange")]):
        y = 1.8 + i * 1.15
        add_arrow(slide, 9.84, y + 0.24, 0.42, 0.2, fill=color)
        add_box(slide, label, 10.35, y, 1.25, 0.62, fill="panel2", line=color, size=9, bold=True)
    add_box(slide, "路由依据：成本 / 授权率 / 地区 / 卡类型 / 故障切换", 6.75, 5.45, 5.1, 0.42, fill="panel", line="navy", size=10, bold=True)


def draw_key_visual(slide):
    add_zone(slide, "密钥分层与信封加密", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    levels = [("Master Key\nHSM / KMS", 8.55, 1.85, 1.8, "red"), ("KEK\n包 DEK", 8.15, 2.8, 2.6, "orange"), ("DEK\n加密 PAN", 7.65, 3.75, 3.6, "cyan")]
    for text, x, y, w, color in levels:
        add_box(slide, text, x, y, w, 0.65, fill="panel2", line=color, size=11, bold=True)
        if y < 3.7:
            add_down_arrow(slide, x + w / 2 - 0.16, y + 0.68, 0.32, 0.28, fill=color)
    add_box(slide, "DB：密文 + 加密后的 DEK", 6.65, 4.8, 2.1, 0.55, fill="panel", line="cyan", size=10, bold=True)
    add_box(slide, "KMS：KEK 与轮换策略", 9.55, 4.8, 2.1, 0.55, fill="panel", line="orange", size=10, bold=True)
    add_arrow(slide, 8.78, 4.98, 0.62, 0.18, fill="navy")


def draw_frontend_visual(slide):
    add_zone(slide, "支付页脚本治理", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "用户浏览器\nPayment Page", 6.55, 2.75, 1.35, 0.9, fill="panel2", line="teal", size=10, bold=True)
    scripts = [("广告", 8.35, 1.75), ("埋点", 9.65, 1.75), ("客服浮窗", 8.35, 4.35), ("A/B 工具", 9.65, 4.35)]
    for label, x, y in scripts:
        add_box(slide, label, x, y, 1.0, 0.45, fill="panel2", line="orange", size=9, bold=True)
    add_box(slide, "CSP / SRI\n清单 / 授权 / 完整性监控", 9.0, 2.75, 2.35, 0.9, fill="navy", line="navy", color="panel", size=10, bold=True)
    add_arrow(slide, 7.9, 3.08, 0.55, 0.2, fill="teal")
    add_arrow(slide, 8.92, 2.2, 0.35, 0.16, fill="orange")
    add_arrow(slide, 8.92, 4.58, 0.35, 0.16, fill="orange")
    add_box(slide, "v4.0 6.4.3：支付页所有脚本必须可证明", 6.8, 5.45, 4.85, 0.42, fill="panel", line="red", size=10, bold=True)


def draw_data_visual(slide):
    add_zone(slide, "数据最小化管道", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "支付事件", 6.55, 3.0, 1.0, 0.5, fill="panel2", line="teal", size=10, bold=True)
    add_arrow(slide, 7.58, 3.15, 0.42, 0.18, fill="teal")
    add_box(slide, "脱敏网关\nPAN→token/hash\nCVV 丢弃", 8.1, 2.65, 1.65, 1.2, fill="navy", line="navy", color="panel", size=10, bold=True)
    sinks = [("日志", 10.15, 1.75, "cyan"), ("客服", 10.15, 2.75, "green"), ("BI", 10.15, 3.75, "orange"), ("AI 风控", 10.15, 4.75, "red")]
    for label, x, y, color in sinks:
        add_arrow(slide, 9.72, y + 0.17, 0.42, 0.16, fill=color)
        add_box(slide, label, x, y, 1.1, 0.42, fill="panel2", line=color, size=9, bold=True)
    add_box(slide, "查找所有副本：ELK / APM / BI / 截图 / 测试环境", 6.75, 5.45, 4.95, 0.42, fill="panel", line="navy", size=10, bold=True)


def draw_devtest_visual(slide):
    add_zone(slide, "研发测试控制链", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    steps = [("代码", "Secrets\n不进 Git", "teal"), ("CI/CD", "变更控制", "cyan"), ("测试", "测试卡\n仿真数据", "orange"), ("生产", "CDE 隔离", "green")]
    for i, (a, b, color) in enumerate(steps):
        x = 6.55 + i * 1.32
        add_box(slide, f"{a}\n{b}", x, 2.25, 1.05, 0.95, fill="panel2", line=color, size=8, bold=True)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.03, 2.58, 0.28, 0.18, fill=color)
    add_box(slide, "SAST / DAST / SCA / ASV / 渗透测试 / 支付页完整性检测", 6.75, 4.25, 4.9, 0.55, fill="navy", line="navy", color="panel", size=10, bold=True)
    add_small_caption(slide, "安全测试结果必须跟踪修复；支付开发需要专项培训", 6.7, 5.35, 5.0)


def draw_incident_visual(slide):
    add_zone(slide, "事件响应时间线", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_arrow(slide, 6.65, 3.05, 4.95, 0.3, fill="navy", text="发现 → 隔离 → 通知 → 取证 → 修复 → 复盘", size=10)
    events = [("0h\n下线/隔离", 6.65, "red"), ("24h\n通知 PSP/收单行", 8.05, "orange"), ("48h+\n取证根因", 9.55, "cyan"), ("演练\n报告保留", 10.85, "green")]
    for text, x, color in events:
        add_box(slide, text, x, 2.05, 1.02, 0.62, fill="panel2", line=color, size=9, bold=True)
        add_down_arrow(slide, x + 0.36, 2.72, 0.28, 0.28, fill=color)
    add_box(slide, "CVV 被记录、密钥泄露都按高危事件处理", 7.0, 4.65, 4.55, 0.55, fill="red", line="red", color="panel", size=11, bold=True)


def draw_evidence_map_visual(slide):
    add_zone(slide, "审计证据地图", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    add_box(slide, "PCI 证据库\n持续收集 / 可追溯", 8.55, 3.0, 1.55, 0.78, fill="navy", line="navy", color="panel", size=10, bold=True)
    items = [
        ("SAQ / AOC", 6.55, 1.85, "teal"),
        ("ROC / QSA", 8.55, 1.85, "cyan"),
        ("ASV 扫描", 10.55, 1.85, "orange"),
        ("渗透测试", 6.55, 4.75, "red"),
        ("变更记录", 8.55, 4.75, "green"),
        ("权限复核", 10.55, 4.75, "navy"),
    ]
    for text, x, y, color in items:
        add_box(slide, text, x, y, 1.1, 0.5, fill="panel2", line=color, size=9, bold=True)
    add_arrow(slide, 7.62, 2.1, 1.0, 0.16, fill="teal")
    add_arrow(slide, 9.2, 2.4, 0.16, 0.48, fill="cyan")
    add_arrow(slide, 10.25, 2.1, 0.45, 0.16, fill="orange")
    add_arrow(slide, 7.62, 4.95, 1.0, 0.16, fill="red")
    add_arrow(slide, 9.2, 4.25, 0.16, 0.48, fill="green")
    add_arrow(slide, 10.25, 4.95, 0.45, 0.16, fill="navy")
    add_box(slide, "v4.0 重点：不是年底补材料，而是持续证明控制有效", 6.85, 5.65, 4.7, 0.42, fill="panel", line="navy", size=10, bold=True)


def draw_closing_visual(slide):
    add_zone(slide, "PCI 业务收益闭环", 6.15, 1.35, 6.25, 5.05, fill="panel", line="navy")
    center = add_box(slide, "降低 Scope\n提升韧性", 8.55, 3.0, 1.45, 0.8, fill="navy", line="navy", color="panel", size=12, bold=True)
    items = [("Network Token\n授权率", 7.0, 1.85, "green"), ("Token 编排\n多 PSP", 10.0, 1.85, "cyan"), ("PayFac\n规模化", 7.0, 4.75, "orange"), ("Passkey / C2P\n少碰卡", 10.0, 4.75, "teal")]
    for text, x, y, color in items:
        add_box(slide, text, x, y, 1.35, 0.62, fill="panel2", line=color, size=9, bold=True)
    add_arrow(slide, 8.35, 2.16, 0.45, 0.18, fill="green")
    add_arrow(slide, 9.76, 2.16, 0.45, 0.18, fill="cyan")
    add_arrow(slide, 8.35, 4.98, 0.45, 0.18, fill="orange")
    add_arrow(slide, 9.76, 4.98, 0.45, 0.18, fill="teal")


VISUALS = {
    "title": draw_title_visual,
    "toc": draw_toc_visual,
    "boundary": draw_boundary_visual,
    "license": draw_license_visual,
    "verification": draw_verification_visual,
    "saq_tree": draw_saq_tree_visual,
    "scope": draw_scope_visual,
    "cde_zoom": draw_cde_zoom_visual,
    "three_lanes": draw_three_lanes_visual,
    "product": draw_product_visual,
    "citmit": draw_citmit_visual,
    "citmit_sequence": draw_citmit_sequence_visual,
    "token": draw_token_visual,
    "network_token_lifecycle": draw_network_token_lifecycle_visual,
    "vault": draw_vault_visual,
    "key": draw_key_visual,
    "frontend": draw_frontend_visual,
    "data": draw_data_visual,
    "devtest": draw_devtest_visual,
    "incident": draw_incident_visual,
    "evidence_map": draw_evidence_map_visual,
    "closing": draw_closing_visual,
}


SLIDES = [
    {
        "title": "PCI DSS 培训要点清单",
        "kicker": "架构图版",
        "bullets": [
            "用产品、架构、运营三种语言解释 PCI DSS",
            "核心问题：是否进 Scope、怎样降成本、出事如何兜底",
            "主线：合规判断 → 产品落地 → 技术架构 → 运营审计",
        ],
        "visual": "title",
        "notes": "开场白：今天不是给大家念合规条文，而是把 PCI DSS 翻译成产品、架构、运营三个人能听懂的语言。四个模块走完，大家应该能回答三件事：我们公司在 Scope 里吗？我们的产品选型把 PCI 成本压到最小了吗？真出事了我们怎么兜底？",
    },
    {
        "title": "目录 · 四大模块",
        "bullets": [
            "1. 合规判断：等级、SAQ、Scope",
            "2. 产品落地：接入方式、CIT/MIT、Token、多 PSP",
            "3. 技术架构：Vault、密钥、前端安全、日志脱敏",
            "4. 运营审计：客服、测试、权限、事件响应",
        ],
        "visual": "toc",
        "notes": "PCI DSS 对很多人来说是年底突击一次的事，但其实它贯穿产品从 0 到 1 的每一个决策。选 Hosted Page 还是 Direct API，Scope 差十倍；要不要自建 Vault，合规成本差百万级。所以今天不讲条文，讲决策。",
    },
    {
        "title": "模块一：合规判断",
        "bullets": ["PCI DSS 是支付责任边界", "先判断 Level / SAQ / Scope", "这三件事答错，后面的产品和架构判断都会跑偏"],
        "module": ("一", ["PCI DSS 是什么", "属于哪一级", "Scope 有多大"]),
        "notes": "模块一解决三个问题：PCI DSS 到底是什么、你属于哪一级、你的 Scope 有多大。这三个问题答错了，后面全错。",
    },
    {
        "title": "补充：PCI 与支付牌照的定位",
        "bullets": [
            "PCI DSS 是卡组网络的安全底线，不是支付经营牌照。",
            "信用卡收单不一定需要卡组监管牌照，但接触、处理、存储持卡人数据就必须满足 PCI DSS。",
            "是否需要支付牌照，取决于当地法律、资金清结算参与方式和业务模式。",
            "先分清监管资质与安全合规，再进入 Level、SAQ、Scope 判断。",
        ],
        "visual": "license",
        "notes": "这一页用来先纠正常见误解：PCI DSS 和支付牌照不是一回事。PCI 是卡组网络的安全准入要求，只要你处理、存储或传输持卡人数据，就必须满足。支付牌照是本地监管和经营许可问题，是否需要取决于具体国家地区、资金清结算方式和业务模式。讲清这个边界以后，再进入认证等级、SAQ 和 Scope，听众会更容易理解为什么 PCI 是底线，而不是牌照。",
    },
    {
        "title": "一、定位与边界",
        "bullets": [
            "PCI DSS 不是安全条文，是支付责任边界：谁接触卡数据、数据流经哪里、哪些系统进入 Scope。",
            "目标不是过审计，是通过产品和架构设计把 PCI 范围降到最小。",
            "v4.0 核心变化：从年度大考变成持续运营。",
        ],
        "visual": "boundary",
        "notes": "第一句是关键：很多团队把 PCI 当 ISO27001 那种通用安全清单来抄，错了。PCI 只看一件事：卡数据碰没碰。你公司有 200 个微服务，可能只有 5 个进 Scope：支付 API、Vault、DB、日志、跳板机。剩下 195 个只要网络隔离做好，跟 PCI 没关系。Hosted Page 还是 Direct API，产品定下来那一刻，未来三年的合规成本就定了。v4.0 的变化是要求持续监控、持续评估，日志告警变更管理都得自动化。",
    },
    {
        "title": "二、认证等级与验证方式",
        "table": [
            ["维度", "要点"],
            ["Merchant Level 1-4", "按年交易量划分，Level 1 >600 万笔"],
            ["SP Level 1-2", "平台 / 网关 / Vault 服务商身份判断"],
            ["SAQ A / A-EP / C / D", "由接入方式决定，不是安全团队决定"],
            ["验证材料", "SAQ、AOC、ROC、ASV 扫描、渗透测试"],
            ["v4.0", "持续合规，非年度 snapshot"],
        ],
        "visual": "verification",
        "notes": "这里最容易混淆的是 SAQ 类型：它是产品形态决定的，不是安全团队选的。Hosted Page 走 SAQ A，iframe/Elements 走 SAQ A-EP，Direct API 走 SAQ D。你让商户从 D 降到 A-EP，合规工作量能少 80%。SP 和 Merchant 的区分也关键：帮自己收单是 Merchant，帮别人收单加存卡是 SP，Level 1 SP 必须 ROC 加 QSA 现场审。",
    },
    {
        "title": "补充：SAQ 决策树",
        "bullets": [
            "SAQ 类型由接入方式决定，不是主观选择。",
            "Hosted Page：卡数据不进入商户系统，通常走 SAQ A。",
            "iframe / Elements：后端不碰卡，但支付页由商户承载，通常走 SAQ A-EP。",
            "Direct API：卡数据进入商户系统，通常走 SAQ D。",
        ],
        "visual": "saq_tree",
        "notes": "这一页把 SAQ 判断讲成决策树。先问卡数据是否进入商户系统。如果完全不进入，Hosted Page 通常是 SAQ A。如果通过 iframe 或 Elements 嵌在商户页面，后端只拿 token，但支付页仍由商户控制，通常是 SAQ A-EP。如果卡号直接进入商户前后端，就是 Direct API，基本走 SAQ D。重点是 SAQ 不是安全团队想选哪个就选哪个，而是产品接入形态的结果。",
    },
    {
        "title": "三、PCI Scope",
        "bullets": [
            "进 Scope 的不只是支付页：API、Vault、DB、缓存、MQ、日志、APM、BI、客服后台、CI/CD、跳板机。",
            "常见盲区：错误日志含 PAN、APM 采请求体、客服截图、测试环境真实卡、风控同步未脱敏数据。",
            "降 Scope 七件套：Hosted Page、iframe、PSP token 替 PAN、CDE 隔离、日志脱敏、最小权限、不存 CVV。",
        ],
        "visual": "scope",
        "notes": "Scope 是 QSA 最爱挑刺的地方。ELK 里一条错误日志打了完整 PAN，整个 ELK 集群进 CDE；APM 采了请求体，APM 进 CDE；客服后台截图里有卡号，客服系统进 CDE。降 Scope 的思路：能不碰就不碰，碰了就脱敏，脱敏不了就隔离，隔离不了就加密。CVV 是铁律：永远不存。",
    },
    {
        "title": "补充：CDE Scope 边界放大图",
        "bullets": [
            "CDE 不只包含支付 API、Vault、PAN DB、KMS/HSM。",
            "日志、APM、客服、BI、MQ、CI/CD 可能因为接触数据或影响 CDE 进入 Scope。",
            "判断标准：是否存储、处理、传输卡数据，或能影响 CDE 安全。",
            "设计目标：缩小 CDE 核心区，把外围系统通过脱敏、隔离和最小权限挡在外面。",
        ],
        "visual": "cde_zoom",
        "notes": "这一页是 Scope 页的放大版。CDE 核心区一般包括支付 API、Vault、PAN 数据库、KMS 或 HSM。但真正容易扩大 Scope 的，是外围系统：日志、APM、客服后台、BI、MQ、缓存、CI/CD。判断标准不是系统名字，而是它是否存储、处理、传输卡数据，或者能影响 CDE 安全。目标是让 CDE 核心区尽量小，外围系统通过脱敏、隔离、权限控制和变更控制挡在外面。",
    },
    {
        "title": "案例一：一条 PAN 日志把 ELK 拉进 Scope",
        "case": {
            "scene": "支付 API 报错时，把完整请求体写入 ELK，里面包含 PAN。",
            "mistake": "团队认为“只有日志系统”，不属于支付核心系统。",
            "impact": "ELK、APM、日志访问权限、留存策略都被 QSA 纳入 Scope。",
            "fix": "入口字段级脱敏；日志禁止 PAN/CVV；APM 关闭请求体采集；定期做数据发现。",
            "takeaway": "Scope 看数据触点，不看系统名字。",
        },
        "notes": "这个案例用来让大家记住 Scope 的判断标准。支付 API 报错时，如果把完整请求体写进 ELK，里面有 PAN，那么日志系统就不再是外围系统，而会被拉进 PCI Scope。误判点是团队以为 ELK 只是运维工具，不是支付系统。正确做法是在入口和日志框架两层做字段级脱敏，禁止 PAN 和 CVV 进入日志，APM 关闭请求体采集，并定期做数据发现。",
    },
    {
        "title": "模块二：产品落地",
        "bullets": ["把合规要求转成产品形态选择", "用 Token、Network Token、多 PSP 编排降低成本", "产品决策决定未来几年的 PCI 成本曲线"],
        "module": ("二", ["接入方式", "CIT / MIT", "Tokenization", "多 PSP 编排"]),
        "notes": "模块一告诉你要求是什么，模块二告诉你怎么设计能把合规成本转成业务收益。产品同学这部分重点听。",
    },
    {
        "title": "补充：资金流 / 数据流 / 责任流",
        "bullets": [
            "同一笔交易里，资金流、卡数据流、合规责任流不是同一条线。",
            "资金可能由 PSP、收单机构、卡组织、发卡行清结算。",
            "卡数据可能只经过卡组织或 PSP，也可能进入商户 CDE。",
            "合规责任按角色和数据触点分配，不按资金路径简单判断。",
        ],
        "visual": "three_lanes",
        "notes": "这一页放在支付产品形态之前，用来建立一个重要认知：钱、数据、责任不是同一条线。资金流讲的是清结算路径，卡数据流讲的是 PAN、token、授权数据经过哪里，责任流讲的是谁因为接触数据或影响 CDE 承担 PCI 控制要求。很多误判来自把这三条线混在一起。比如商户不直接清算资金，不代表不承担 PCI；商户使用 PSP，也不代表所有系统都自动出 Scope。",
    },
    {
        "title": "四、支付产品形态与 PCI 影响",
        "table": [
            ["形态", "PCI 范围", "适用"],
            ["Hosted Payment Page", "最小（SAQ A）", "跳转可接受"],
            ["Embedded / Elements", "小（SAQ A-EP）", "前端担责，后端只拿 token"],
            ["Direct API", "最大（SAQ D）", "自建风控 / 多 PSP 路由 / 自建 Vault"],
            ["PayFac", "平台扛 Level 1", "子商户快速 onboarding"],
            ["Click to Pay / Passkey", "近零", "卡组织 / FIDO2 托管"],
        ],
        "visual": "product",
        "notes": "四种形态对应三种 SAQ，合规成本差一个数量级。Hosted Page 最省但跳转伤转化；Elements 是甜点区，前端要管 CSP/SRI；Direct API 除非有强理由，否则别碰。PayFac 平台做 Master Merchant，子商户不用单独收单，但平台要扛 PCI Level 1。Click to Pay 和 Passkey 的好处是卡数据存在卡组织或用户设备，商户连碰都碰不到。",
    },
    {
        "title": "案例二：Direct API 让商户掉到 SAQ D",
        "case": {
            "scene": "为了减少一次跳转，产品把 Hosted Page 改成自建卡输入页 + Direct API。",
            "mistake": "只评估转化率，没有评估 SAQ 和 CDE 成本变化。",
            "impact": "商户从 SAQ A / A-EP 进入 SAQ D，前后端、日志、客服、CI/CD 都要重做控制。",
            "fix": "优先使用 Elements 或 Click to Pay；确需 Direct API 时，先做 CDE 设计和成本评审。",
            "takeaway": "支付体验优化不能绕开 PCI 成本评估。",
        },
        "notes": "这个案例讲给产品同学。为了减少跳转，自建卡输入页再走 Direct API，看起来只是体验优化，但它可能让商户从 SAQ A 或 A-EP 掉到 SAQ D。影响不是多填几张表，而是前端安全、后端隔离、日志脱敏、客服权限、CI/CD 变更控制全部进入更重的控制范围。正确做法是优先评估 Elements、Click to Pay 或 Hosted Page 的体验优化空间。如果确实必须 Direct API，要先完成 CDE 设计和合规成本评审。",
    },
    {
        "title": "补充：Passkey + Click to Pay 现代支付架构",
        "full_image": "passkey_click_to_pay_architecture.png",
        "notes": "这一页补模块二的关键认知：Passkey 和 Click to Pay 不是同一层能力。Passkey 解决身份认证，用 WebAuthn/FIDO2 让用户无密码登录，降低账号接管和钓鱼风险。Click to Pay 解决支付凭证和支付发起，卡号以网络 token 或卡组侧凭证存在卡组织或发卡行体系中，商户前端只发起支付请求，尽量不接触卡号。两者组合以后，用户体验是无密码登录加一键支付，合规效果是商户不存密码、不存卡号、不处理卡数据，PCI Scope 大幅收缩。这页也要强调：Passkey 不是支付本身，Click to Pay 也不是登录本身，它们分别解决身份和支付凭证问题。",
    },
    {
        "title": "五、CIT / MIT 与业务场景",
        "bullets": [
            "CIT = 买家在场：首笔、3DS、SCA。",
            "MIT = 商户发起：订阅续费、免密、分期。",
            "绑卡：存 token 不存卡号，首笔必 CIT + 3DS。",
            "订阅 = CIT 首笔 + MIT 后续，需维护授权凭证、重试、卡过期处理。",
        ],
        "visual": "citmit",
        "notes": "CIT/MIT 是订阅、复购、一键支付的地基。首次必须 CIT + 3DS，建立授权；后续 MIT 必须带首次凭证，金额频率不能超授权范围。订阅业务最怕两件事：扣款失败重试没做好、卡过期没更新。Network Token 能解决卡过期自动换新 token。",
    },
    {
        "title": "补充：3DS + CIT/MIT 时序图",
        "bullets": [
            "首笔 CIT：用户在场，完成 3DS / SCA，建立授权。",
            "PSP / 发卡行返回授权结果，商户保存 token 或 agreement。",
            "后续 MIT：商户基于首次授权凭证发起续费或免密扣款。",
            "关键约束：金额、频率、场景不能超出用户授权范围。",
        ],
        "visual": "citmit_sequence",
        "notes": "这一页用时序图讲订阅的合规逻辑。第一步是 CIT 首笔，用户在场，完成 3DS 或 SCA。第二步是 PSP 和发卡行返回授权结果。第三步是商户保存 token 或 agreement，注意不是保存卡号。第四步是后续 MIT，商户基于首次授权凭证发起续费或免密扣款。重点提醒：MIT 不是任意扣款，它必须引用首次凭证，金额、频率、场景不能超出用户授权范围。",
    },
    {
        "title": "六、Tokenization",
        "table": [
            ["Token 类型", "跨 PSP", "可逆", "备注"],
            ["PSP token", "否", "通常由 PSP 可逆", "换 PSP 失效"],
            ["Merchant token", "是", "是（经 Vault）", "自建合规成本高"],
            ["Network token", "强", "否", "卡组织发，授权率 +2~5%"],
            ["Payment method ID", "否", "否", "绑定用户级"],
        ],
        "visual": "token",
        "notes": "Token 是产品落地最核心的一页。四个维度判断：可逆性、跨 PSP、泄露风险、绑定粒度。Network Token 不只是安全，也是业务增长引擎：授权率提升、卡过期自动更新、绑定商户和设备。多 PSP 场景的痛点也在这里：PSP A 的 token 在 PSP B 用不了，所以多 PSP 策略往往倒逼独立 Vault 或 Network Token。",
    },
    {
        "title": "补充：Network Token 生命周期",
        "bullets": [
            "商户或 PSP 向卡组织 token 服务请求 Network Token。",
            "卡组织发放 token，并绑定商户、设备或使用场景。",
            "授权时用 token 替代 PAN，降低卡号暴露面。",
            "卡过期、换卡、补卡时 token 可自动更新。",
            "多 PSP 场景下，Network Token 能降低 PSP token 锁定问题。",
        ],
        "visual": "network_token_lifecycle",
        "notes": "这一页补 Network Token 的生命周期。先由商户或 PSP 请求 token，卡组织发放 Network Token，并绑定商户、设备或使用场景。授权时用 token 替代 PAN，降低卡号暴露面。卡过期、换卡、补卡时，卡组织可以自动更新 token，减少订阅扣款失败。多 PSP 场景下，Network Token 能缓解 PSP token 被某一家锁死的问题，但具体跨 PSP 能力仍取决于接入方式、地区和卡组织支持。",
    },
    {
        "title": "案例三：Network Token 降低订阅扣款失败",
        "case": {
            "scene": "订阅业务大量失败来自卡过期、换卡、补卡后凭证失效。",
            "mistake": "只把 token 当脱敏工具，没有利用卡组织 token 生命周期能力。",
            "impact": "续费失败、被动流失增加，多 PSP 路由也受 PSP token 锁定限制。",
            "fix": "优先接入 Network Token；结合账户更新、重试策略和多 PSP 路由。",
            "takeaway": "合规设计也可以转成授权率和续费率收益。",
        },
        "notes": "这个案例用来说明 PCI 不只是成本。订阅业务里，大量扣款失败来自卡过期、换卡、补卡以后原凭证失效。如果只把 token 当脱敏工具，就会错过 Network Token 的生命周期能力。Network Token 可以在卡生命周期变化时自动更新，减少续费失败；在多 PSP 场景下，也能缓解单一 PSP token 锁定的问题。这里的结论是：好的合规设计可以转化成授权率、续费率和容灾能力。",
    },
    {
        "title": "模块三：技术架构",
        "bullets": ["Vault 怎么选", "密钥怎么管", "前端支付页怎么防 Magecart", "日志、客服、BI、AI 风控怎么做数据最小化"],
        "module": ("三", ["Card Vault", "密钥管理", "前端支付页安全", "数据治理"]),
        "notes": "模块三给架构和工程同学。Vault 怎么选、密钥怎么管、前端支付页怎么防 Magecart、日志怎么脱敏。",
    },
    {
        "title": "七、Card Vault 选型",
        "table": [
            ["方案", "成本", "跨 PSP", "适合"],
            ["PSP Vault", "最低", "否", "单 PSP"],
            ["第三方 Vault", "中", "是", "多 PSP，不自建"],
            ["自建 Vault", "最高", "是", "大体量 + 强控制权"],
        ],
        "bullets": ["多 PSP 场景下，独立 Vault 的价值在编排：成本路由、授权率路由、容灾切换、卡类型路由。"],
        "visual": "vault",
        "notes": "Vault 选型看三件事：PSP 数量、交易量、团队能力。单 PSP 就用 PSP Vault；3 家以上 PSP 就要考虑独立 Vault 做编排。自建 Vault 的合规成本一年几十万到上百万起步，HSM/KMS、双人控制、渗透测试、审计，没专门团队别碰。第三方 Vault 是中间路线，选之前一定要看对方 AOC。",
    },
    {
        "title": "八、密钥管理",
        "bullets": [
            "分层模型：DEK → KEK → Master Key，每层只加密下一层。",
            "信封加密：密文 + 加密后的 DEK 存 DB，KEK 在 KMS，密钥与密文分离。",
            "HSM / KMS：加密 PAN 的密钥必须在 HSM 或等效设备中。",
            "生命周期：生成 → 分发 → 使用 → 存储 → 轮换 → 吊销 → 销毁。",
            "v4.0：密钥必须定期轮换，轮换策略必须有文档。",
        ],
        "visual": "key",
        "notes": "密钥管理的常见翻车点：硬编码进代码、密钥和密文同库、生产密钥拷到测试环境、邮件发密钥。分层和信封加密的重点是即使数据库被拖库，攻击者拿到的也是密文和被 KEK 包过的 DEK，没有 KEK 解不开。KEK 在 KMS/HSM，拖库没用。",
    },
    {
        "title": "九、前端支付页安全",
        "bullets": [
            "威胁：Magecart 恶意 JS 窃卡。英航 2018，22 行代码跑 22 天，38 万卡号泄露。",
            "风险来源：广告脚本、埋点 SDK、客服浮窗、A/B 工具、CDN 篡改。",
            "控制：CSP + SRI + 脚本白名单 + 变更审批 + 完整性监控。",
            "v4.0 6.4.3：支付页所有脚本必须有清单、有授权、有完整性监控。",
        ],
        "visual": "frontend",
        "notes": "前端这页很多人忽略。卡数据不走服务器也不代表安全，Magecart 直接在用户浏览器里偷。英航案例可以强调：22 行恶意代码，CDN 被篡改，22 天才发现。v4.0 的 6.4.2 和 6.4.3 是硬要求：支付页每个脚本都要有 inventory、authorization、integrity monitoring。",
    },
    {
        "title": "十、数据治理（日志 / 客服 / BI / AI 风控）",
        "bullets": [
            "日志：全链路不记录 PAN/CVV，正则脱敏，权限 + 留存周期。",
            "客服：只展示后四位 + 卡品牌，全量审计，权限拆分。",
            "BI/数仓：PAN → token/hash，CVV 不进，BIN 可进。",
            "AI 风控 vs 数据最小化：用 BIN / token / PSP 风控分 / 合成数据平衡。",
            "v4.0：数据发现分类、留存处置策略、访问权限定期复核。",
        ],
        "visual": "data",
        "notes": "数据治理讲卡数据除了 CDE 里那份，其他副本在哪。ELK、APM trace、BI、客服截图、测试环境，每份都要追踪。AI 风控和数据最小化有冲突，解法包括 BIN 代 PAN、token 代用户键、用 PSP 风控分、合成数据训练。",
    },
    {
        "title": "模块四：运营审计",
        "bullets": ["研发测试怎么管", "审计材料怎么备", "出事以后怎么响应", "持续证明控制有效，而不是年底突击"],
        "module": ("四", ["研发与测试", "审计材料", "事件响应"]),
        "notes": "最后一个模块给运营、运维、合规同学。研发测试怎么管、审计材料怎么备、出事了怎么响应。",
    },
    {
        "title": "补充：审计证据地图",
        "bullets": [
            "审计不是年底补材料，而是日常证据持续沉淀。",
            "核心材料：SAQ、AOC、ROC、ASV 扫描、渗透测试。",
            "运营证据：变更记录、权限复核、安全测试修复闭环、事件演练记录。",
            "目标：能证明控制持续有效，而不是只证明某一天有效。",
        ],
        "visual": "evidence_map",
        "notes": "这一页讲审计证据地图。PCI v4.0 的重点不是年底补材料，而是持续证明控制有效。核心材料包括 SAQ、AOC、ROC、ASV 扫描和渗透测试。运营证据包括变更记录、权限复核、安全测试修复闭环、事件响应演练记录。审计时最怕的是控制写在制度里，但没有执行证据。所以证据库要日常沉淀，最好和工单、CI/CD、扫描平台、权限系统打通。",
    },
    {
        "title": "十一、研发与测试",
        "bullets": [
            "测试环境：禁真实卡，用 PSP 测试卡或脱敏仿真数据。",
            "密钥：不进 Git，Secrets Manager，进了立即轮换。",
            "CI/CD：支付代码更严变更控制，影响 CDE 的变更须正式流程。",
            "安全测试：SAST + DAST + SCA + ASV 扫描 + 渗透测试 + 支付页完整性检测。",
            "v4.0：开发/测试/生产隔离，安全测试结果须跟踪修复，支付开发须专项培训。",
        ],
        "visual": "devtest",
        "notes": "研发这页的红线两条：测试环境真实卡、密钥进 Git。v4.0 新增的几条：三环境隔离、安全测试结果不能测完就扔、支付开发要有专项安全培训。",
    },
    {
        "title": "十二、审计与事件响应",
        "table": [
            ["场景", "关键动作"],
            ["PAN 泄露", "24h 内通知 PSP/收单行/卡组织，取证 + 根因"],
            ["CVV 被记录", "更严重，本身违规，立即删 + 全面排查"],
            ["密钥泄露", "立即轮换 + 重加密 + 等同 PAN 泄露处置"],
            ["支付页 JS 篡改", "下线页面 + 取证 + 清理 + 加强 CSP/SRI 后上线"],
            ["v4.0", "事件响应计划每年演练，关键人 24/7 可达"],
        ],
        "visual": "incident",
        "notes": "事件响应这页，时间线最关键：PAN 泄露多数卡组织要求 24 小时内通知，错过窗口罚金翻倍。CVV 被记录比 PAN 泄露更严重，因为存储 CVV 本身就是 PCI 禁止的。密钥泄露是最高级事件，等同于 PAN 泄露。v4.0 要求事件响应计划每年至少演练一次，关键人 24/7 联系得上，复盘报告保留三年。",
    },
    {
        "title": "培训主线 · 收尾",
        "bullets": [
            "PCI DSS 不只是把卡数据暴露面降到最小。",
            "通过 Network Token、Token 编排、PayFac、Click to Pay / Passkey，把合规成本转化为业务收益。",
            "业务收益：授权率提升、churn 降低、多 PSP 容灾、商户规模化覆盖。",
            "四个模块逻辑：要求是什么 → 怎么设计降成本 → 怎么落地 → 怎么证明做到了。",
        ],
        "visual": "closing",
        "notes": "收尾把今天的线串一下：PCI 很多人当合规负担，但它其实能转成业务收益。Network Token 提授权率，Token 编排撑多 PSP 路由，PayFac 撑商户规模化，Click to Pay/Passkey 降 Scope 还提转化。四个模块的逻辑链：合规判断告诉你要求是什么，产品落地告诉你怎么设计降成本，技术架构告诉你怎么落地，运营审计告诉你怎么持续证明做到了。",
    },
]


def build_slide(prs, spec, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if "full_image" in spec:
        draw_full_image(slide, spec["full_image"])
        add_text(slide, f"{index:02d}", 12.25, 6.92, 0.45, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)
        slide.notes_slide.notes_text_frame.text = spec.get("notes", "")
        return slide

    add_header(slide, spec["title"], spec.get("kicker"))
    if "case" in spec:
        draw_case_card(slide, spec["case"])
        add_text(slide, f"{index:02d}", 12.25, 6.92, 0.45, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)
        slide.notes_slide.notes_text_frame.text = spec.get("notes", "")
        return slide

    if "table" in spec:
        add_table(slide, spec["table"], 0.65, 1.55, 5.0, 3.1, font_size=8.5)
        if spec.get("bullets"):
            add_bullets(slide, spec["bullets"], x=0.68, y=4.85, w=5.0, h=1.0, size=12)
    else:
        add_bullets(slide, spec.get("bullets", []), x=0.65, y=1.55, w=5.05, h=4.75, size=14)

    if "module" in spec:
        module_no, labels = spec["module"]
        draw_module_visual(slide, module_no, labels)
    elif "visual" in spec:
        VISUALS[spec["visual"]](slide)

    add_text(slide, f"{index:02d}", 12.25, 6.92, 0.45, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)
    slide.notes_slide.notes_text_frame.text = spec.get("notes", "")
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, spec in enumerate(SLIDES, start=1):
        build_slide(prs, spec, idx)
    prs.save(OUT)
    print(f"generated {OUT} with {len(SLIDES)} slides")


if __name__ == "__main__":
    main()
